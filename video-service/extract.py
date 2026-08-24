import fitz
from pptx import Presentation
from pathlib import Path


def extract_document(path: str):
    """Return a list of slide/page dictionaries from PPTX, PDF, or TXT."""
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pptx":
        prs = Presentation(path)
        pages = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
            pages.append({"number": i, "text": "\n".join(parts).strip()})
        return pages

    if suffix == ".pdf":
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc, 1):
            pages.append({"number": i, "text": page.get_text("text").strip()})
        doc.close()
        return pages

    if suffix == ".txt":
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
        chunks = [x.strip() for x in text.split("\f") if x.strip()]
        if not chunks:
            chunks = [text.strip()]
        return [{"number": i, "text": chunk} for i, chunk in enumerate(chunks, 1)]

    raise ValueError("Only PPTX, PDF and TXT files are supported.")
